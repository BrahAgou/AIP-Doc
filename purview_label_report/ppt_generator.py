from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_ppt(labels, output_file):
    prs = Presentation()
    # Custom title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Purview Sensitivity Labels Report"
    slide.placeholders[1].text = "Exported from Microsoft Purview via Graph API"
    
    for label in labels:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = label['label_name']
        # Advanced styling for the content
        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        tf.clear()
        p = tf.add_paragraph()
        p.text = f"Description: {label.get('description', '')}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 51, 102)
        tf.add_paragraph()
        p = tf.add_paragraph()
        p.text = "Publishing Policy:"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(44, 117, 255)
        for policy in label['policies']:
            pol_p = tf.add_paragraph()
            pol_p.text = f"• {policy['type']}: {policy['name']} ({policy['user_count']} users)"
            pol_p.level = 1
            pol_p.font.size = Pt(12)
            pol_p.font.color.rgb = RGBColor(90, 90, 90)
            pol_p.alignment = PP_ALIGN.LEFT
    prs.save(output_file)